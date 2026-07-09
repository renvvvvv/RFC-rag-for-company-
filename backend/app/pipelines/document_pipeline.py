import os
import re
from typing import List, Dict, Any, Tuple
from uuid import UUID
from app.pipelines.base import BaseIngestPipeline


class DocumentIngestPipeline(BaseIngestPipeline):
    """处理 doc/docx/pdf/pptx/txt/markdown/html/json 文档。

    支持：
    - PDF 文本抽取，失败页面使用 OCR 回退
    - 表格抽取并包裹 [TABLE]...[/TABLE] 标记
    - 元数据提取（title/author/page_count/language）
    - 外部依赖缺失时优雅降级并返回 TODO 占位提示
    """

    @property
    def supported_types(self) -> List[str]:
        return [
            "document", "pdf", "txt", "markdown", "md",
            "html", "htm", "pptx", "ppt", "json",
            "docx", "doc",
        ]

    def process(
        self,
        file_path: str,
        doc_id: UUID,
        metadata: Dict[str, Any] = None,
    ) -> List[Dict[str, Any]]:
        metadata = metadata or {}
        ext = os.path.splitext(file_path)[1].lower()

        extracted_text = ""
        doc_metadata: Dict[str, Any] = {}

        if ext in [".docx", ".doc"]:
            extracted_text, doc_metadata = self._extract_docx(file_path)
        elif ext == ".pdf":
            extracted_text, doc_metadata = self._extract_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            extracted_text, doc_metadata = self._extract_pptx(file_path)
        elif ext in [".txt", ".md", ".markdown"]:
            extracted_text, doc_metadata = self._extract_txt(file_path)
        elif ext in [".html", ".htm"]:
            extracted_text, doc_metadata = self._extract_html(file_path)
        elif ext == ".json":
            extracted_text, doc_metadata = self._extract_json(file_path)
        else:
            extracted_text = ""

        merged_meta = {**metadata, **doc_metadata}
        # 2026-07-08: 文本清洗（在切分前执行）
        extracted_text = self._clean_text(extracted_text)
        chunks = self._chunk_text(extracted_text)
        for chunk in chunks:
            chunk["metadata"] = {**chunk.get("metadata", {}), **merged_meta}
        return chunks

    def _extract_docx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        try:
            import docx
            doc = docx.Document(file_path)
            parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)

            # 表格抽取：仅做标记，不参与布局顺序还原
            for table in doc.tables:
                parts.append("[TABLE]")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    parts.append(row_text)
                parts.append("[/TABLE]")

            text = "\n".join(parts)
            core_props = doc.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "page_count": None,  # DOCX 无固定页数
                "language": self._detect_language(text),
            }
            return text, metadata
        except Exception as e:
            return f"[解析失败: {str(e)}]", {}

    def _extract_pdf(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        try:
            import pdfplumber
            texts = []
            metadata: Dict[str, Any] = {}

            with pdfplumber.open(file_path) as pdf:
                if pdf.metadata:
                    metadata = {
                        "title": pdf.metadata.get("Title", ""),
                        "author": pdf.metadata.get("Author", ""),
                    }
                metadata["page_count"] = len(pdf.pages)

                for page_idx, page in enumerate(pdf.pages):
                    text = page.extract_text()

                    if not text or not text.strip():
                        # OCR 回退：扫描型 PDF 页面
                        ocr_text = self._ocr_pdf_page(file_path, page_idx)
                        if ocr_text and ocr_text != "[OCR_UNAVAILABLE]":
                            texts.append(f"[OCR_PAGE {page_idx + 1}]\n{ocr_text}\n[/OCR_PAGE]")
                    else:
                        # 表格标记
                        page_tables = page.extract_tables() or []
                        for table in page_tables:
                            if table:
                                texts.append("[TABLE]")
                                for row in table:
                                    texts.append(
                                        " | ".join(
                                            str(cell) if cell is not None else ""
                                            for cell in row
                                        )
                                    )
                                texts.append("[/TABLE]")
                        texts.append(text)

            metadata["language"] = self._detect_language("\n".join(texts))
            return "\n".join(texts), metadata
        except Exception as e:
            return f"[解析失败: {str(e)}]", {}

    def _ocr_pdf_page(self, file_path: str, page_number: int) -> str:
        """使用 pdf2image + pytesseract 对单页 PDF 做 OCR。"""
        try:
            from pdf2image import convert_from_path
            import pytesseract

            images = convert_from_path(
                file_path,
                first_page=page_number + 1,
                last_page=page_number + 1,
            )
            if images:
                return pytesseract.image_to_string(images[0], lang="chi_sim+eng")
            return ""
        except Exception:
            # TODO: 安装 pdf2image 与 tesseract 后启用 PDF OCR 回退。
            return "[OCR_UNAVAILABLE]"

    def _extract_pptx(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        """解析 PPT/PPTX 文件，支持 .ppt（旧格式）和 .pptx（新格式）"""
        try:
            import subprocess
            import tempfile
            from pathlib import Path

            actual_path = file_path

            # 2026-07-06: 支持 .ppt 旧格式
            # 2026-07-08: 优化 - 先复制到简单路径避免中文路径问题
            if file_path.lower().endswith('.ppt') and not file_path.lower().endswith('.pptx'):
                import shutil
                tmp_input = '/tmp/convert_input.ppt'
                tmp_output = '/tmp/convert_input.pptx'

                try:
                    # 复制到简单路径（避免中文路径问题）
                    shutil.copy2(file_path, tmp_input)

                    # 转换成 pptx
                    result = subprocess.run(
                        ['libreoffice', '--headless', '--convert-to', 'pptx', tmp_input, '--outdir', '/tmp/'],
                        capture_output=True, text=True, timeout=120
                    )

                    if Path(tmp_output).exists():
                        actual_path = tmp_output
                    else:
                        return f"[PPT转换失败: 输出文件不存在]", {}
                except Exception as e:
                    return f"[PPT转换失败: {str(e)}]", {}
                finally:
                    # 清理临时文件
                    Path(tmp_input).unlink(missing_ok=True)

            # 用 python-pptx 解析
            from pptx import Presentation
            prs = Presentation(actual_path)
            texts = []

            for slide_idx, slide in enumerate(prs.slides):
                slide_texts = []
                slide_texts.append(f"--- 第{slide_idx +1}页 ---")

                for shape in slide.shapes:
                    try:
                        # 提取普通文本
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_texts.append(shape.text.strip())

                        # 提取表格内容
                        if shape.has_table:
                            table = shape.table
                            for row in table.rows:
                                row_texts = []
                                for cell in row.cells:
                                    if cell.text.strip():
                                        row_texts.append(cell.text.strip())
                                if row_texts:
                                    slide_texts.append(" | ".join(row_texts))
                    except Exception:
                        # 跳过无法读取的形状
                        continue

                if len(slide_texts) > 1:  # 有实际内容
                    texts.append("\n".join(slide_texts))

            core_props = prs.core_properties
            metadata = {
                "title": core_props.title or "",
                "author": core_props.author or "",
                "page_count": len(prs.slides),
                "language": self._detect_language("\n".join(texts)),
            }

            # 清理转换后的临时文件
            if actual_path == '/tmp/convert_input.pptx':
                Path(actual_path).unlink(missing_ok=True)

            return "\n\n".join(texts), metadata
        except Exception as e:
            return f"[解析失败: {str(e)}]", {}

    def _extract_txt(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        metadata = {
            "title": "",
            "author": "",
            "page_count": 1,
            "language": self._detect_language(text),
        }
        return text, metadata

    def _extract_html(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        try:
            from bs4 import BeautifulSoup
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            text = soup.get_text(separator="\n")
            title = soup.title.string.strip() if soup.title and soup.title.string else ""
            metadata = {
                "title": title,
                "author": "",
                "page_count": 1,
                "language": self._detect_language(text),
            }
            return text, metadata
        except Exception as e:
            return f"[解析失败: {str(e)}]", {}

    def _extract_json(self, file_path: str) -> Tuple[str, Dict[str, Any]]:
        try:
            import json
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
            metadata = {
                "title": "",
                "author": "",
                "page_count": 1,
                "language": self._detect_language(text),
            }
            return text, metadata
        except Exception as e:
            return f"[解析失败: {str(e)}]", {}

    def _detect_language(self, text: str) -> str:
        if not text or not text.strip():
            return ""
        try:
            from langdetect import detect
            return detect(text[:2000])
        except Exception:
            # TODO: 安装 langdetect 后启用语言检测。
            return ""

    # 2026-07-03: 原来的切分逻辑（按固定字符数切分），已注释保留
    # def _chunk_text(
    #     self,
    #     text: str,
    #     chunk_size: int = 500,
    #     overlap: int = 100,
    # ) -> List[Dict[str, Any]]:
    #     chunks = []
    #     if not text:
    #         return chunks
    #
    #     start = 0
    #     chunk_index = 0
    #     while start < len(text):
    #         end = min(start + chunk_size, len(text))
    #         chunk_text = text[start:end]
    #
    #         chunks.append({
    #             "content": chunk_text,
    #             "modality": "text",
    #             "chunk_index": chunk_index,
    #             "position_info": {"paragraph_range": [start, end]},
    #             "metadata": {},
    #         })
    #
    #         start += chunk_size - overlap
    #         chunk_index += 1
    #
    #     return chunks

    # 2026-07-03: 优化的切分逻辑（按段落切分 + 合并小chunks + 拆分大chunks + 保持overlap）
    # 2026-07-03: 调整参数 - min_size=500, target_size=600, max_size=800, overlap=100
    def _chunk_text(
        self,
        text: str,
        target_size: int = 600,
        min_size: int = 500,
        max_size: int = 800,
        overlap: int = 100,
    ) -> List[Dict[str, Any]]:
        """优化的切分方案：按段落切分，保持语义完整性"""
        chunks = []
        if not text:
            return chunks

        # 1. 先按段落切分（遇到空行就切）
        paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]

        # 如果段落数量太少，尝试按单个换行符切分
        if len(paragraphs) < 3:
            paragraphs = [p.strip() for p in text.split('\n') if p.strip()]

        max_para_len = max((len(p) for p in paragraphs), default=0)

        # 2. 合并小段落
        merged = []
        current = ""
        for para in paragraphs:
            if len(current) + len(para) <= target_size:
                current += "\n\n" + para if current else para
            else:
                if current:
                    merged.append(current)
                current = para
        if current:
            merged.append(current)


        # 3. 拆分大段落（特殊处理表格和目录）
        for chunk in merged:
            # 判断是否是表格块（以[TABLE]开头[/TABLE]结尾）、目录（不拆分）
            stripped = chunk.strip()
            is_table = stripped.startswith('[TABLE]') and stripped.endswith('[/TABLE]')
            # 目录块：大部分行包含 '.....'  点划线（且不是正文中偶然出现的省略号）
            dir_lines = [l for l in stripped.split('\n') if '.....' in l]
            is_directory = len(dir_lines) >= 5 and len(dir_lines) > len(stripped.split('\n')) * 0.5
            if is_table:
                # 表格块不拆分
                chunks.append(chunk)
            elif is_directory and len(chunk) <= max_size:
                # 小目录块不拆分
                chunks.append(chunk)
            elif len(chunk) > max_size or (is_directory and len(chunk) > max_size):
                # 非表格内容才拆分
                sentences = self._split_by_sentence(chunk)
                sub_chunks = self._merge_sentences(sentences, target_size)

                # 2026-07-08: 兜底逻辑 — 句子切分后仍有超大块（>2000字符）的按 \n 硬切
                # 注意：只对非表格块生效（表格块以 [TABLE] 开头、[/TABLE] 结尾）
                final_sub = []
                for sub in sub_chunks:
                    is_table_block = sub.strip().startswith('[TABLE]') and sub.strip().endswith('[/TABLE]')
                    if len(sub) > 2000 and not is_table_block:
                        lines = sub.split('\n')
                        current_line = ""
                        for line in lines:
                            if len(current_line) + len(line) <= target_size:
                                current_line += '\n' + line if current_line else line
                            else:
                                if current_line:
                                    final_sub.append(current_line)
                                current_line = line
                        if current_line:
                            final_sub.append(current_line)
                    else:
                        final_sub.append(sub)
                sub_chunks = final_sub

                for sub_chunk in sub_chunks:
                    chunks.append(sub_chunk)
            else:
                chunks.append(chunk)


        # 4. 强制合并小 chunks
        chunks = self._merge_small_chunks_aggressive(chunks, min_size)

        # 5. 添加 overlap
        chunks_with_overlap = []
        for i, chunk in enumerate(chunks):
            if i > 0:
                # 添加上一个 chunk 的最后 overlap 个字符
                prev = chunks[i-1]
                overlap_text = prev[-overlap:] if len(prev) > overlap else prev
                chunk = overlap_text + "\n\n" + chunk
            chunks_with_overlap.append(chunk)

        # 5. 格式化输出
        result = []
        for i, chunk_text in enumerate(chunks_with_overlap):
            result.append({
                "content": chunk_text,
                "modality": "text",
                "chunk_index": i,
                "position_info": {"chunk_index": i},
                "metadata": {},
            })

        return result

    # def _merge_small_chunks_aggressive(self, chunks: List[str], min_size: int) -> List[str]:
        """强制合并小 chunks，确保每个 chunk 至少 min_size 字符"""
        max_before = max((len(c) for c in chunks), default=0)
        result = self.__merge_aggressive(chunks, min_size)
        max_after = max((len(r) for r in result), default=0)
        return result

    def _merge_small_chunks_aggressive(self, chunks: List[str], min_size: int) -> List[str]:
        if not chunks:
            return chunks

        merged = []
        current = chunks[0]

        for chunk in chunks[1:]:
            if len(current) < min_size:
                # 当前 chunk 太小，强制合并
                current += "\n\n" + chunk
            else:
                merged.append(current)
                current = chunk

        merged.append(current)

        # 再次检查，确保没有太小的 chunks
        final_merged = []
        for chunk in merged:
            if len(chunk) < min_size and final_merged:
                # 合并到前一个 chunk
                final_merged[-1] += "\n\n" + chunk
            else:
                final_merged.append(chunk)

        return final_merged

    def _split_by_sentence(self, text: str) -> List[str]:
        """按句子切分（遇到句号就切）"""
        import re
        # 匹配中英文句号、感叹号、问号
        pattern = r'(?<=[。！？.!?])\s*'
        return [s.strip() for s in re.split(pattern, text) if s.strip()]

    def _merge_sentences(self, sentences: List[str], target_size: int) -> List[str]:
        """合并句子到合适的 chunk 大小"""
        chunks = []
        current_chunk = ""

        for sentence in sentences:
            if len(current_chunk) + len(sentence) <= target_size:
                current_chunk += sentence
            else:
                if current_chunk:
                    chunks.append(current_chunk)
                current_chunk = sentence

        if current_chunk:
            chunks.append(current_chunk)

        return chunks
